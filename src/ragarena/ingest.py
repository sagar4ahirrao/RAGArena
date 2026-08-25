"""Multi-format document ingestion for RagArena.

Supports the most common RAG source materials and degrades gracefully when an
optional dependency is missing.  Returns a list of document dicts compatible
with :class:`ragarena.index.VectorIndex` / :func:`ragarena.engine.evaluate`:

    {"text": str, "metadata": {...}, "tables": [[...]], "images": [path_or_b64]}
"""
from __future__ import annotations

import base64
import csv
import json
import os
import re
from typing import Any, Dict, List, Optional

try:
    from .index import MultimodalDocument  # type: ignore
except Exception:  # pragma: no cover
    MultimodalDocument = None  # type: ignore


# ──────────────────────────────────────────────────────────────────────────────
# Public API
# ──────────────────────────────────────────────────────────────────────────────
def parse_file(path: str, parser: str = "auto") -> List[Dict[str, Any]]:
    """Parse *path* into one or more document dicts.

    parser: "auto" (by extension) | "text" | "markdown" | "html" | "raw"
    """
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    meta = {"source": os.path.basename(path), "path": path, "type": ext}

    if ext in ("txt", "md", "markdown"):
        return _text(path, meta, markdown=(ext in ("md", "markdown")))
    if ext == "pdf":
        return _pdf(path, meta, parser)
    if ext == "docx":
        return _docx(path, meta)
    if ext == "pptx":
        return _pptx(path, meta)
    if ext in ("html", "htm"):
        return _html(path, meta)
    if ext == "csv":
        return _csv(path, meta)
    if ext == "tsv":
        return _csv(path, meta, delimiter="\t")
    if ext in ("json", "jsonl"):
        return _json(path, meta)
    if ext == "xml":
        return _xml(path, meta)
    if ext in ("yaml", "yml"):
        return _yaml(path, meta)
    if ext in ("xlsx", "xls"):
        return _xlsx(path, meta)
    if ext in ("png", "jpg", "jpeg"):
        return _image(path, meta)
    if ext == "sql":
        return _sql_file(path, meta)
    if ext in ("db", "sqlite", "sqlite3"):
        return _sqlite_db(path, meta)
    # Fallback: treat as plain text
    return _text(path, meta, markdown=False)


def parse_dir(root: str, recursive: bool = True, parser: str = "auto") -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    files = []
    for dirpath, _, fnames in os.walk(root):
        for f in fnames:
            files.append(os.path.join(dirpath, f))
        if not recursive:
            break
    for fp in sorted(files):
        try:
            out.extend(parse_file(fp, parser))
        except Exception as e:  # skip unreadable files
            out.append({"text": "", "metadata": {"source": fp, "error": str(e)[:120]}})
    return [d for d in out if d.get("text") or d.get("tables") or d.get("images")]


# ──────────────────────────────────────────────────────────────────────────────
# Format handlers
# ──────────────────────────────────────────────────────────────────────────────
def _read(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _text(path: str, meta: dict, markdown: bool) -> List[Dict[str, Any]]:
    text = _read(path)
    if markdown:
        try:
            import markdown as md_module
            html = md_module.markdown(text)
            text = re.sub(r"<[^>]+>", " ", html)
            text = re.sub(r"\s+", " ", text).strip()
        except Exception:
            pass
    return [{"text": text, "metadata": {**meta, "parser": "markdown" if markdown else "text"}}]


def _pdf(path: str, meta: dict, parser: str) -> List[Dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except Exception:
        return [{"text": _read(path), "metadata": {**meta, "parser": "raw"}}]
    reader = PdfReader(path)
    docs: List[Dict[str, Any]] = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        doc = {"text": txt.strip(),
               "metadata": {**meta, "page": i + 1, "parser": "pdf"}}
        try:  # tables via pdfplumber (optional)
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                tbl = pdf.pages[i].extract_tables()
                if tbl:
                    doc["tables"] = [[ [c for c in row] for row in t] for t in tbl]
        except Exception:
            pass
        if doc["text"] or doc.get("tables"):
            docs.append(doc)
    return docs or [{"text": "", "metadata": meta}]


def _docx(path: str, meta: dict) -> List[Dict[str, Any]]:
    from docx import Document
    doc = Document(path)
    parts: List[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = [[ [c.text for c in r.cells] for r in t.rows] for t in doc.tables]
    text = "\n\n".join(parts)
    return [{"text": text, "tables": tables,
             "metadata": {**meta, "parser": "docx"}}]


def _pptx(path: str, meta: dict) -> List[Dict[str, Any]]:
    from pptx import Presentation
    prs = Presentation(path)
    out: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        parts: List[str] = []
        tables: List[List[List[str]]] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                tables.append([[c.text for c in row.cells] for row in shape.table.rows])
        out.append({"text": "\n".join(parts), "tables": tables,
                    "metadata": {**meta, "slide": i + 1, "parser": "pptx"}})
    return [d for d in out if d["text"] or d["tables"]]


def _html(path: str, meta: dict) -> List[Dict[str, Any]]:
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(_read(path), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = re.sub(r"\s+", " ", soup.get_text(" ")).strip()
    tables: List[List[List[str]]] = []
    for t in soup.find_all("table"):
        rows = [[c.get_text(strip=True) for c in r.find_all(["td", "th"])]
                for r in t.find_all("tr")]
        if rows:
            tables.append(rows)
    return [{"text": text, "tables": tables,
             "metadata": {**meta, "parser": "html"}}]


def _flatten(obj: Any, prefix: str = "", max_depth: int = 12) -> List[str]:
    """Flatten nested dicts/lists into "a.b.c: value" lines — every line is a
    self-contained retrievable fact with its full key-path as context."""
    out: List[str] = []
    if max_depth <= 0:
        return [f"{prefix}: {obj}" if prefix else str(obj)]
    if isinstance(obj, dict):
        for k, v in obj.items():
            out.extend(_flatten(v, f"{prefix}.{k}" if prefix else str(k), max_depth - 1))
    elif isinstance(obj, list):
        if obj and all(not isinstance(v, (dict, list)) for v in obj):
            out.append(f"{prefix}: {', '.join(str(v) for v in obj)}")
        else:
            for i, v in enumerate(obj):
                out.extend(_flatten(v, f"{prefix}[{i}]" if prefix else f"[{i}]", max_depth - 1))
    else:
        out.append(f"{prefix}: {obj}" if prefix else str(obj))
    return out


def _row_groups(header: List[Any], rows: List[List[Any]], meta: dict,
                group_size: int = 50, parser: str = "csv") -> List[Dict[str, Any]]:
    """Chunk table rows into documents of *group_size* rows. Every chunk repeats
    the column headers so each embedded piece is self-contained ("col: value")."""
    hdr = [str(c) for c in header]
    docs: List[Dict[str, Any]] = []
    total = 0
    for g in range(0, len(rows), group_size):
        grp = rows[g:g + group_size]
        text = "\n".join(
            ", ".join(f"{hdr[i]}: {c}" for i, c in enumerate(r) if i < len(hdr))
            for r in grp if r
        )
        if not text.strip():
            continue
        docs.append({
            "text": text,
            "tables": [[hdr] + [[str(c) for c in r] for r in grp]],
            "metadata": {**meta, "parser": parser,
                         "rows": f"{g + 1}-{g + len(grp)}", "n_rows": len(grp)},
        })
        total += len(grp)
    if total == 0 and hdr and any(h.strip() for h in hdr):
        docs.append({"text": ", ".join(hdr), "tables": [[hdr]],
                     "metadata": {**meta, "parser": parser, "rows": "header-only"}})
    return docs


def _csv(path: str, meta: dict, delimiter: str = ",") -> List[Dict[str, Any]]:
    rows: List[List[str]] = []
    with open(path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        reader = csv.reader(f, delimiter=delimiter)
        for r in reader:
            rows.append(r)
    if not rows:
        return []
    docs = _row_groups(rows[0], rows[1:], meta,
                       parser="tsv" if delimiter == "\t" else "csv")
    if docs:
        docs[-1].setdefault("tables", [])
        docs[-1]["tables"] = [[[str(c) for c in r] for r in rows]]
    return docs


def _xml(path: str, meta: dict) -> List[Dict[str, Any]]:
    """Generic XML: walk the tree emitting 'path/to/tag: text' plus attributes."""
    from xml.etree import ElementTree as ET
    try:
        root = ET.fromstring(_read(path))
    except Exception:
        return [{"text": _read(path), "metadata": {**meta, "parser": "xml-raw"}}]
    lines: List[str] = []

    def walk(el: ET.Element, prefix: str, depth: int = 0) -> None:
        tag = f"{prefix}/{el.tag}" if prefix else el.tag
        for k, v in el.attrib.items():
            lines.append(f"{tag}@{k}: {v}")
        txt = (el.text or "").strip()
        if txt:
            lines.append(f"{tag}: {txt}")
        if depth < 20:
            for child in el:
                walk(child, tag, depth + 1)

    walk(root, "")
    return [{"text": "\n".join(lines), "metadata": {**meta, "parser": "xml"}}]


def _yaml(path: str, meta: dict) -> List[Dict[str, Any]]:
    import yaml as yaml_module
    try:
        data = yaml_module.safe_load(_read(path))
    except Exception:
        return [{"text": _read(path), "metadata": {**meta, "parser": "yaml-raw"}}]
    if isinstance(data, list):
        return [{"text": "\n".join(_flatten(item)),
                 "metadata": {**meta, "parser": "yaml", "record": i}}
                for i, item in enumerate(data[:500])]
    text = "\n".join(_flatten(data or {}))
    return [{"text": text, "metadata": {**meta, "parser": "yaml"}}]


def _json(path: str, meta: dict) -> List[Dict[str, Any]]:
    raw = _read(path)
    try:
        data = json.loads(raw)
    except Exception:
        # maybe JSON-Lines
        items = []
        for line in raw.splitlines():
            line = line.strip().rstrip(",")
            if not line:
                continue
            try:
                items.append(json.loads(line))
            except Exception:
                pass
        if items:
            data = items
        else:
            return [{"text": raw, "metadata": {**meta, "parser": "json-raw"}}]

    if isinstance(data, list):
        cap = 500
        docs = [{"text": "\n".join(_flatten(item)),
                 "metadata": {**meta, "parser": "json", "record": i}}
                for i, item in enumerate(data[:cap])]
        if len(data) > cap:
            docs.append({"text": f"[{len(data) - cap} more records truncated]",
                         "metadata": {**meta, "parser": "json", "record": "truncated"}})
        return docs

    return [{"text": "\n".join(_flatten(data)), "metadata": {**meta, "parser": "json"}}]


def _xlsx(path: str, meta: dict) -> List[Dict[str, Any]]:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    out: List[Dict[str, Any]] = []
    for ws in wb.worksheets:
        rows = [[str(c.value) if c.value is not None else "" for c in r] for r in ws.iter_rows()]
        if not rows:
            continue
        sheet_meta = {**meta, "sheet": ws.title}
        out.extend(_row_groups(rows[0], rows[1:], sheet_meta, parser="xlsx"))
    return out


def _image(path: str, meta: dict) -> List[Dict[str, Any]]:
    b64 = ""
    try:
        with open(path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode()
    except Exception:
        pass
    return [{"text": "", "images": [b64],
             "metadata": {**meta, "parser": "image"}}]


def _sql_file(path: str, meta: dict) -> List[Dict[str, Any]]:
    """A .sql dump/script: split into individual statements, one document per statement."""
    raw = _read(path)
    statements = [s.strip() for s in raw.split(";") if s.strip()]
    return [{"text": s, "metadata": {**meta, "parser": "sql", "statement_index": i}}
            for i, s in enumerate(statements)] or [{"text": raw, "metadata": {**meta, "parser": "sql"}}]


def _sqlite_db(path: str, meta: dict, max_rows_per_table: int = 500) -> List[Dict[str, Any]]:
    """A SQLite database file: one document per table, rows rendered as readable text + a raw table."""
    import sqlite3
    docs: List[Dict[str, Any]] = []
    conn = sqlite3.connect(path)
    try:
        cur = conn.cursor()
        cur.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%'")
        tables = [r[0] for r in cur.fetchall()]
        for table in tables:
            cur.execute(f'SELECT * FROM "{table}" LIMIT {max_rows_per_table}')
            cols = [d[0] for d in cur.description]
            rows = cur.fetchall()
            text = "\n".join(
                ", ".join(f"{c}: {v}" for c, v in zip(cols, row)) for row in rows
            )
            docs.append({
                "text": text,
                "tables": [[cols] + [list(r) for r in rows]],
                "metadata": {**meta, "parser": "sqlite", "table": table, "n_rows": len(rows)},
            })
    finally:
        conn.close()
    return docs


def from_sql(connection_string: str, query: str, source_name: str = "sql") -> List[Dict[str, Any]]:
    """Pull rows from ANY SQL database (Postgres/MySQL/SQLite/...) via SQLAlchemy and
    turn each row into a retrievable document. Requires ``pip install ragarena[sql]``
    plus the relevant DB driver (e.g. ``psycopg2-binary``, ``pymysql``).

    Example::

        docs = from_sql("postgresql://user:pass@host/db",
                         "SELECT id, title, body FROM articles")
    """
    from sqlalchemy import create_engine, text as sql_text
    engine = create_engine(connection_string)
    with engine.connect() as conn:
        result = conn.execute(sql_text(query))
        cols = list(result.keys())
        rows = [list(r) for r in result.fetchall()]
    return [{
        "text": ", ".join(f"{c}: {v}" for c, v in zip(cols, row)),
        "tables": [[cols] + rows] if rows else None,
        "metadata": {"source": source_name, "type": "sql", "parser": "sql_query", "row_index": i},
    } for i, row in enumerate(rows)]


def to_multimodal(docs: List[Dict[str, Any]]) -> List[Any]:
    """Expand parsed docs into MultimodalDocument objects (text/table/image) when possible."""
    if MultimodalDocument is None:
        return docs
    out: List[Any] = []
    for d in docs:
        meta = d.get("metadata", {})
        if d.get("text"):
            out.append(MultimodalDocument(content=d["text"], doc_type="text", metadata=meta))
        for i, tbl in enumerate(d.get("tables") or []):
            rendered = "\n".join(", ".join(str(c) for c in row) for row in tbl)
            out.append(MultimodalDocument(content=rendered, doc_type="table",
                                           metadata={**meta, "table_index": i}))
        for i, img in enumerate(d.get("images") or []):
            out.append(MultimodalDocument(content=img, doc_type="image",
                                           metadata={**meta, "image_index": i}))
    return out
